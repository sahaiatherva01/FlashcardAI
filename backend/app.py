"""
app.py
-------------------------------------------------------------------
AI Flashcard Generator — Backend (B.Tech Mini Project edition)

A single-file Flask backend that:
  1. Serves the static frontend (frontend/index.html, style.css, script.js)
  2. Accepts up to 10 files (PDF/DOCX/TXT/PPTX/image) or pasted text
  3. Extracts + cleans + de-duplicates + chunks the text
  4. Sends each chunk to Gemini with an engineered prompt
  5. Validates the returned JSON flashcards (rejecting duplicates too)
     and stores them in SQLite
  6. Exposes REST endpoints for listing/editing/deleting/favoriting/
     quizzing/exporting, plus an upload-history endpoint

Kept intentionally as ONE file since this is a mini/academic project
rather than a production system.
-------------------------------------------------------------------
"""

import os
import re
import csv
import json
import uuid
import random
import sqlite3
import logging
from io import StringIO, BytesIO
from datetime import datetime, timezone

from flask import Flask, request, jsonify, send_from_directory, g, Response, send_file
from werkzeug.utils import secure_filename
from dotenv import load_dotenv

import fitz                     # PyMuPDF -> PDF text
import docx                     # python-docx -> DOCX text
from pptx import Presentation   # python-pptx -> PPTX text
from PIL import Image, UnidentifiedImageError
import pytesseract              # OCR for images
from google import genai                # Gemini (Google GenAI SDK)
from google.genai import types
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib import colors

# ============================================================
# CONFIG
# ============================================================

load_dotenv()

BASE_DIR = os.path.abspath(os.path.dirname(__file__))
FRONTEND_DIR = os.path.join(BASE_DIR, "..", "frontend")
DB_PATH = os.path.join(BASE_DIR, "flashcards.db")
UPLOAD_DIR = os.path.join(BASE_DIR, "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
GEMINI_MODEL = os.environ.get("GEMINI_MODEL", "gemini-2.5-flash")
ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "pptx", "png", "jpg", "jpeg", "webp"}
CARD_TYPES = ("definition", "qa", "fill_blank", "mcq", "formula", "true_false")
QUIZ_TYPES = ("mcq", "true_false", "fill_blank")
MAX_FILES_PER_UPLOAD = 10

logging.basicConfig(level=logging.INFO, format="[%(asctime)s] %(levelname)s: %(message)s")
logger = logging.getLogger("flashcard_app")

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 20 * 1024 * 1024 * MAX_FILES_PER_UPLOAD  # generous ceiling; per-file checked separately


# ============================================================
# DATABASE (plain sqlite3 — no ORM needed for a project this size)
# ============================================================

def get_db():
    if "db" not in g:
        g.db = sqlite3.connect(DB_PATH)
        g.db.row_factory = sqlite3.Row
    return g.db


@app.teardown_appcontext
def close_db(exception=None):
    db = g.pop("db", None)
    if db is not None:
        db.close()


def _add_column_if_missing(conn, table, column, ddl):
    """Lightweight migration helper — lets older DBs upgrade in place."""
    existing = {row[1] for row in conn.execute(f"PRAGMA table_info({table})")}
    if column not in existing:
        conn.execute(f"ALTER TABLE {table} ADD COLUMN {ddl}")


def init_db():
    conn = sqlite3.connect(DB_PATH)

    conn.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            filename TEXT NOT NULL,
            source_type TEXT NOT NULL,
            subject TEXT,
            status TEXT NOT NULL DEFAULT 'processing',
            error_message TEXT,
            card_count INTEGER DEFAULT 0,
            created_at TEXT NOT NULL
        )
    """)

    conn.execute("""
        CREATE TABLE IF NOT EXISTS flashcards (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            card_type TEXT NOT NULL,
            subject TEXT,
            topic TEXT,
            difficulty TEXT DEFAULT 'medium',
            term TEXT,
            question TEXT,
            answer TEXT NOT NULL,
            explanation TEXT,
            options TEXT,
            document_id INTEGER,
            favorite INTEGER DEFAULT 0,
            created_at TEXT NOT NULL
        )
    """)

    # Migrations for anyone upgrading from the previous minimal schema.
    _add_column_if_missing(conn, "flashcards", "document_id", "document_id INTEGER")
    _add_column_if_missing(conn, "flashcards", "favorite", "favorite INTEGER DEFAULT 0")

    conn.commit()
    conn.close()


# ============================================================
# TEXT EXTRACTION (PDF / DOCX / TXT / PPTX / Image OCR)
# ============================================================

class ExtractionError(Exception):
    """Raised with a user-friendly message when a file can't be read."""


def extract_text(filepath: str, ext: str) -> str:
    """Dispatch to the right extractor based on file extension. Never
    fails silently — every failure mode raises ExtractionError with a
    message that can be shown directly to the user."""
    try:
        if ext == "pdf":
            with fitz.open(filepath) as pdf:
                text = "\n".join(page.get_text() for page in pdf).strip()
            if not text:
                raise ExtractionError(
                    "No selectable text found — this PDF may be scanned. "
                    "Try uploading it as an image instead so OCR can run."
                )
            return text

        if ext == "docx":
            document = docx.Document(filepath)
            text = "\n".join(p.text for p in document.paragraphs if p.text.strip())
            if not text:
                raise ExtractionError("This DOCX file appears to be empty.")
            return text

        if ext == "pptx":
            presentation = Presentation(filepath)
            parts = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text.strip())
            text = "\n".join(parts)
            if not text:
                raise ExtractionError("No text found on any slide in this PPTX.")
            return text

        if ext == "txt":
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read().strip()
            if not text:
                raise ExtractionError("This text file is empty.")
            return text

        if ext in {"png", "jpg", "jpeg", "webp"}:
            try:
                image = Image.open(filepath).convert("L")  # grayscale improves OCR
            except UnidentifiedImageError:
                raise ExtractionError("This image file is corrupted or unreadable.")
            text = pytesseract.image_to_string(image).strip()
            if not text:
                raise ExtractionError(
                    "OCR could not find readable text in this image. Try a clearer photo."
                )
            return text

        raise ExtractionError(f"Unsupported file type: .{ext}")

    except ExtractionError:
        raise
    except fitz.FileDataError:
        raise ExtractionError("This PDF file is corrupted and could not be opened.")
    except Exception as exc:
        logger.exception("Unexpected extraction failure for %s", filepath)
        raise ExtractionError(f"Could not read this file: {exc}")


def clean_text(text: str) -> str:
    """Basic normalization before chunking/prompting."""
    text = text.replace("\r\n", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def dedupe_paragraphs(texts: list) -> str:
    """
    Merges multiple documents' text into one, dropping paragraphs that
    are exact duplicates (common when the same slide/section appears in
    more than one uploaded file). Order-preserving, case-insensitive match.
    """
    seen = set()
    merged_paragraphs = []
    for text in texts:
        for para in re.split(r"\n\s*\n", text):
            para = para.strip()
            if not para:
                continue
            key = re.sub(r"\s+", " ", para.lower())
            if key in seen:
                continue
            seen.add(key)
            merged_paragraphs.append(para)
    return "\n\n".join(merged_paragraphs)


def chunk_text(text: str, max_chars: int = 2200) -> list:
    """
    Splits text into roughly max_chars-sized chunks on sentence boundaries
    (a lightweight regex-based splitter — no heavy NLP dependency needed
    for a project this size).
    """
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks, current = [], ""
    for sentence in sentences:
        if len(current) + len(sentence) > max_chars and current:
            chunks.append(current.strip())
            current = sentence
        else:
            current += " " + sentence
    if current.strip():
        chunks.append(current.strip())
    return chunks


# ============================================================
# PROMPT ENGINEERING + GEMINI
# ============================================================

def build_prompt(chunk: str, subject: str | None) -> str:
    """
    Engineered prompt: closed card-type list, strict JSON schema,
    explicit anti-hallucination rule, difficulty criteria, and an
    explicit instruction to avoid duplicate/repeated questions.
    """
    subject_line = f'Subject: "{subject}".' if subject else ""
    return f"""You are an expert study assistant creating exam-ready flashcards for a
university student. Only use facts explicitly present in the STUDY TEXT below —
never invent definitions, numbers, or examples.

{subject_line}

Generate 4-8 flashcards using ONLY these types: definition, qa, fill_blank, mcq, formula, true_false.
- "formula" cards ONLY if an actual formula/equation appears in the text.
- "mcq" cards need exactly 4 options, with "answer" matching one option exactly.
- Assign difficulty: "easy" (single stated fact), "medium" (connects two ideas), "hard" (applies a formula/process).
- Do NOT repeat the same fact in two different cards — each card must test a distinct concept.
- If the text has nothing educational, return [].

Return ONLY a valid JSON array, no markdown fences, in this shape:
[
  {{"card_type": "definition", "term": "...", "question": null, "answer": "...", "explanation": null, "options": [], "topic": "...", "difficulty": "easy"}},
  {{"card_type": "mcq", "term": null, "question": "...", "answer": "...", "explanation": "...", "options": ["...","...","...","..."], "topic": "...", "difficulty": "medium"}}
]

STUDY TEXT:
\"\"\"{chunk}\"\"\"
"""


_gemini_client = None


def _get_gemini_client():
    """Lazily builds a single reusable Gemini client (avoids reconnecting
    on every chunk of every upload)."""
    global _gemini_client
    if _gemini_client is None:
        _gemini_client = genai.Client(
            api_key=GEMINI_API_KEY,
            http_options=types.HttpOptions(timeout=30_000),  # 30s, ms
        )
    return _gemini_client


def call_gemini(prompt: str):
    """Calls Gemini and parses the JSON response, stripping markdown fences if present."""
    if not GEMINI_API_KEY:
        raise RuntimeError("GEMINI_API_KEY is not set. Add it to backend/.env")

    client = _get_gemini_client()
    response = client.models.generate_content(
        model=GEMINI_MODEL,
        contents=prompt,
        config=types.GenerateContentConfig(temperature=0.4, response_mime_type="application/json"),
    )
    raw = (response.text or "").strip()
    if raw.startswith("```"):
        raw = raw.strip("`")
        if raw.lower().startswith("json"):
            raw = raw[4:]
    return json.loads(raw.strip())


# ============================================================
# OFFLINE FALLBACK — used only if Gemini is unreachable/unset
# ============================================================

_DEFINITION_PATTERN = re.compile(
    r"^(?P<term>[A-Z][A-Za-z0-9 \-/]{2,60}?)\s+(?:is|are|refers to|means|denotes)\s+(?P<definition>.{15,300})$"
)


def generate_basic_flashcards(text: str, max_cards: int = 15) -> list:
    """
    Lightweight, non-AI flashcard generator used as a safety net when Gemini
    can't be reached (missing/invalid key, quota exceeded, network error).
    It never raises — worst case it returns an empty list — so the app can
    always tell the user *something* instead of hard-failing the whole upload.

    Strategy: scan sentences for a "Term is/are/means ..." pattern to build
    definition cards, and turn other substantial sentences into
    fill-in-the-blank cards on their last key word/phrase.
    """
    cards = []
    seen_answers = set()

    for paragraph in re.split(r"\n\s*\n", text):
        for sentence in re.split(r"(?<=[.!?])\s+", paragraph.strip()):
            sentence = sentence.strip()
            if len(sentence) < 25 or len(sentence) > 320:
                continue

            match = _DEFINITION_PATTERN.match(sentence)
            if match and len(match.group("term").split()) <= 6:
                term = match.group("term").strip()
                definition = match.group("definition").strip().rstrip(".")
                key = term.lower()
                if key in seen_answers:
                    continue
                seen_answers.add(key)
                cards.append({
                    "card_type": "definition", "term": term, "question": None,
                    "answer": definition, "explanation": None, "options": [],
                    "topic": None, "difficulty": "easy",
                })
            else:
                words = sentence.rstrip(".").split()
                if len(words) < 7:
                    continue
                blank_word = words[-1].strip(",.;:")
                if len(blank_word) < 3 or blank_word.lower() in seen_answers:
                    continue
                seen_answers.add(blank_word.lower())
                blanked = " ".join(words[:-1]) + " ____"
                cards.append({
                    "card_type": "fill_blank", "term": None, "question": blanked,
                    "answer": blank_word, "explanation": None, "options": [],
                    "topic": None, "difficulty": "medium",
                })

            if len(cards) >= max_cards:
                return cards
    return cards


# ============================================================
# VALIDATION — never trust AI output blindly before saving to DB
# ============================================================

def validate_card(card: dict) -> dict | None:
    if not isinstance(card, dict):
        return None
    card_type = card.get("card_type")
    answer = (card.get("answer") or "").strip()
    if card_type not in CARD_TYPES or not answer:
        return None

    if card_type == "mcq":
        options = card.get("options") or []
        if not isinstance(options, list) or len(options) < 2 or answer not in options:
            return None

    if card_type in ("qa", "fill_blank", "mcq", "true_false") and not (card.get("question") or "").strip():
        return None
    if card_type in ("definition", "formula") and not (card.get("term") or "").strip():
        return None

    return {
        "card_type": card_type,
        "term": (card.get("term") or "").strip() or None,
        "question": (card.get("question") or "").strip() or None,
        "answer": answer,
        "explanation": (card.get("explanation") or "").strip() or None,
        "options": card.get("options") or [],
        "topic": (card.get("topic") or "").strip()[:150] or None,
        "difficulty": card.get("difficulty") if card.get("difficulty") in ("easy", "medium", "hard") else "medium",
    }


def _card_fingerprint(card: dict) -> str:
    """Normalized signature used to detect duplicate/near-duplicate cards."""
    prompt_text = (card["term"] or card["question"] or "").strip().lower()
    prompt_text = re.sub(r"\s+", " ", prompt_text)
    return f"{card['card_type']}::{prompt_text}"


def _save_cards(db, raw_cards, subject, document_id, existing_fingerprints) -> int:
    """Validates + dedupes + inserts a batch of raw card dicts. Returns count saved."""
    saved = 0
    for raw_card in raw_cards:
        card = validate_card(raw_card)
        if not card:
            logger.warning("Rejected malformed card: %s", raw_card)
            continue

        fingerprint = _card_fingerprint(card)
        if fingerprint in existing_fingerprints:
            logger.info("Skipped duplicate card: %s", fingerprint[:80])
            continue
        existing_fingerprints.add(fingerprint)

        db.execute(
            """INSERT INTO flashcards
               (card_type, subject, topic, difficulty, term, question, answer,
                explanation, options, document_id, favorite, created_at)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, 0, ?)""",
            (card["card_type"], subject, card["topic"], card["difficulty"], card["term"],
             card["question"], card["answer"], card["explanation"], json.dumps(card["options"]),
             document_id, datetime.now(timezone.utc).isoformat()),
        )
        saved += 1
    return saved


def generate_and_store(raw_text: str, subject: str | None, document_id: int | None = None) -> tuple[int, bool]:
    """Runs clean -> chunk -> prompt -> Gemini -> validate -> dedupe -> store.
    Returns (cards_saved, used_fallback). If Gemini fails on every single
    chunk (missing/invalid key, quota exceeded, network error, etc.) this
    falls back to a lightweight offline generator instead of returning
    nothing — the app should never completely fail an upload."""
    text = clean_text(raw_text)
    chunks = chunk_text(text)
    if not chunks:
        raise ValueError("No usable text found to generate flashcards from.")

    db = get_db()

    # Pre-load fingerprints of existing cards for this subject so we don't
    # save near-duplicates of flashcards the student already has.
    existing_fingerprints = set()
    if subject:
        rows = db.execute("SELECT card_type, term, question FROM flashcards WHERE subject = ?", (subject,)).fetchall()
        for r in rows:
            existing_fingerprints.add(_card_fingerprint(dict(r)))

    saved = 0
    gemini_failures = 0
    for chunk in chunks:
        try:
            raw_cards = call_gemini(build_prompt(chunk, subject))
        except Exception as exc:
            logger.error("Gemini call failed for a chunk: %s", exc)
            gemini_failures += 1
            continue

        if not isinstance(raw_cards, list):
            logger.warning("Gemini returned a non-list response; skipping chunk.")
            continue

        saved += _save_cards(db, raw_cards, subject, document_id, existing_fingerprints)

    used_fallback = False
    if saved == 0 and gemini_failures == len(chunks):
        logger.warning("Gemini unavailable for every chunk — using offline fallback generator.")
        used_fallback = True
        saved = _save_cards(db, generate_basic_flashcards(text), subject, document_id, existing_fingerprints)

    db.commit()
    return saved, used_fallback


def row_to_dict(row: sqlite3.Row) -> dict:
    d = dict(row)
    if "options" in d:
        d["options"] = json.loads(d["options"]) if d.get("options") else []
    if "favorite" in d:
        d["favorite"] = bool(d["favorite"])
    return d


# ============================================================
# ROUTES — Frontend
# ============================================================

@app.route("/")
def serve_index():
    return send_from_directory(FRONTEND_DIR, "index.html")


@app.route("/<path:filename>")
def serve_static(filename):
    return send_from_directory(FRONTEND_DIR, filename)


# ============================================================
# ROUTES — Upload (multi-file, up to 10) + Paste
# ============================================================

@app.route("/api/upload", methods=["POST"])
def api_upload():
    """
    Accepts up to MAX_FILES_PER_UPLOAD files under the "files" field.
    mode="combined" (default): merge all extracted text (de-duplicated by
        paragraph) into ONE flashcard deck.
    mode="separate": generate an independent deck per file.
    A failure on one file never stops processing of the others.
    """
    files = request.files.getlist("files") or (
        [request.files["file"]] if "file" in request.files else []
    )
    subject = (request.form.get("subject") or "").strip() or None
    mode = request.form.get("mode", "combined")

    if not files or all(f.filename == "" for f in files):
        return jsonify({"error": "No files provided."}), 400

    if len(files) > MAX_FILES_PER_UPLOAD:
        return jsonify({"error": f"Please upload at most {MAX_FILES_PER_UPLOAD} files at a time."}), 400

    db = get_db()
    file_results = []       # per-file status for the response
    extracted_texts = []    # (document_id, text) for successfully extracted files

    for file in files:
        if not file or file.filename == "":
            continue

        filename = secure_filename(file.filename)
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

        if ext not in ALLOWED_EXTENSIONS:
            file_results.append({"filename": filename, "status": "failed", "error": "Unsupported file type."})
            continue

        # Prefix with a short unique id so two uploads with the same
        # filename (e.g. two different students both naming a file
        # "notes.pdf") never collide or overwrite each other on disk.
        unique_name = f"{uuid.uuid4().hex[:8]}_{filename}"
        save_path = os.path.join(UPLOAD_DIR, unique_name)
        file.save(save_path)

        cursor = db.execute(
            "INSERT INTO documents (filename, source_type, subject, status, created_at) VALUES (?, ?, ?, 'processing', ?)",
            (filename, ext, subject, datetime.now(timezone.utc).isoformat()),
        )
        document_id = cursor.lastrowid
        db.commit()

        try:
            text = extract_text(save_path, ext)
            extracted_texts.append((document_id, text))
            file_results.append({"filename": filename, "status": "extracted", "document_id": document_id})
        except ExtractionError as exc:
            db.execute("UPDATE documents SET status='failed', error_message=? WHERE id=?", (str(exc), document_id))
            db.commit()
            file_results.append({"filename": filename, "status": "failed", "error": str(exc)})
            logger.warning("Extraction failed for %s: %s", filename, exc)
        finally:
            if os.path.exists(save_path):
                os.remove(save_path)

    if not extracted_texts:
        return jsonify({
            "error": "No files could be processed.",
            "results": file_results,
        }), 422

    total_generated = 0
    any_fallback = False
    try:
        if mode == "separate":
            for document_id, text in extracted_texts:
                count, used_fallback = generate_and_store(text, subject, document_id=document_id)
                any_fallback = any_fallback or used_fallback
                db.execute("UPDATE documents SET status='ready', card_count=? WHERE id=?", (count, document_id))
                total_generated += count
                for r in file_results:
                    if r.get("document_id") == document_id:
                        r["status"] = "ready"
                        r["card_count"] = count
        else:  # combined
            merged_text = dedupe_paragraphs([t for _, t in extracted_texts])
            primary_document_id = extracted_texts[0][0]
            total_generated, any_fallback = generate_and_store(merged_text, subject, document_id=primary_document_id)
            for document_id, _ in extracted_texts:
                db.execute("UPDATE documents SET status='ready', card_count=? WHERE id=?",
                           (total_generated if document_id == primary_document_id else 0, document_id))
            for r in file_results:
                if r.get("status") == "extracted":
                    r["status"] = "ready"
        db.commit()
    except Exception as exc:
        logger.exception("Flashcard generation failed")
        return jsonify({"error": f"Flashcard generation failed: {exc}", "results": file_results}), 500

    succeeded = sum(1 for r in file_results if r["status"] == "ready")
    failed = sum(1 for r in file_results if r["status"] == "failed")
    fallback_note = " (Gemini was unavailable, so basic flashcards were generated instead.)" if any_fallback else ""

    return jsonify({
        "message": f"Processed {len(file_results)} file(s): {succeeded} succeeded, {failed} failed. "
                   f"Generated {total_generated} flashcards.{fallback_note}",
        "count": total_generated,
        "used_fallback": any_fallback,
        "results": file_results,
    })


@app.route("/api/paste", methods=["POST"])
def api_paste():
    """Handles pasted text -> generate flashcards."""
    data = request.get_json(force=True) or {}
    text = (data.get("text") or "").strip()
    subject = (data.get("subject") or "").strip() or None

    if len(text) < 40:
        return jsonify({"error": "Please paste at least a few sentences."}), 400

    db = get_db()
    cursor = db.execute(
        "INSERT INTO documents (filename, source_type, subject, status, created_at) VALUES (?, 'paste', ?, 'processing', ?)",
        ("Pasted text", subject, datetime.now(timezone.utc).isoformat()),
    )
    document_id = cursor.lastrowid
    db.commit()

    try:
        count, used_fallback = generate_and_store(text, subject, document_id=document_id)
        db.execute("UPDATE documents SET status='ready', card_count=? WHERE id=?", (count, document_id))
        db.commit()
        note = " (Gemini was unavailable, so basic flashcards were generated instead.)" if used_fallback else ""
        return jsonify({"message": f"Generated {count} flashcards.{note}", "count": count, "used_fallback": used_fallback})
    except Exception as exc:
        db.execute("UPDATE documents SET status='failed', error_message=? WHERE id=?", (str(exc), document_id))
        db.commit()
        logger.exception("Paste pipeline failed")
        return jsonify({"error": str(exc)}), 500


# ============================================================
# ROUTES — Flashcards CRUD, Favorites, Search/Filter
# ============================================================

@app.route("/api/flashcards", methods=["GET"])
def api_list_flashcards():
    """Lists flashcards with optional search/filter query params."""
    search = request.args.get("q", "").strip()
    subject = request.args.get("subject", "").strip()
    card_type = request.args.get("card_type", "").strip()
    difficulty = request.args.get("difficulty", "").strip()
    favorite_only = request.args.get("favorite", "").strip() == "true"

    query = "SELECT * FROM flashcards WHERE 1=1"
    params = []
    if search:
        query += " AND (term LIKE ? OR question LIKE ? OR answer LIKE ? OR topic LIKE ?)"
        like = f"%{search}%"
        params += [like, like, like, like]
    if subject:
        query += " AND subject = ?"
        params.append(subject)
    if card_type:
        query += " AND card_type = ?"
        params.append(card_type)
    if difficulty:
        query += " AND difficulty = ?"
        params.append(difficulty)
    if favorite_only:
        query += " AND favorite = 1"
    query += " ORDER BY created_at DESC"

    try:
        rows = get_db().execute(query, params).fetchall()
    except sqlite3.Error as exc:
        logger.exception("Database error while listing flashcards")
        return jsonify({"error": f"Database error: {exc}"}), 500

    return jsonify([row_to_dict(r) for r in rows])


@app.route("/api/flashcards/<int:card_id>", methods=["PUT"])
def api_update_flashcard(card_id):
    data = request.get_json(force=True) or {}
    db = get_db()
    row = db.execute("SELECT * FROM flashcards WHERE id = ?", (card_id,)).fetchone()
    if not row:
        return jsonify({"error": "Flashcard not found."}), 404

    fields = {k: data.get(k, row[k]) for k in ("term", "question", "answer", "explanation", "topic", "difficulty")}
    db.execute(
        """UPDATE flashcards SET term=?, question=?, answer=?, explanation=?, topic=?, difficulty=? WHERE id=?""",
        (fields["term"], fields["question"], fields["answer"], fields["explanation"],
         fields["topic"], fields["difficulty"], card_id),
    )
    db.commit()
    return jsonify({"message": "Flashcard updated."})


@app.route("/api/flashcards/<int:card_id>", methods=["DELETE"])
def api_delete_flashcard(card_id):
    db = get_db()
    row = db.execute("SELECT id FROM flashcards WHERE id = ?", (card_id,)).fetchone()
    if not row:
        return jsonify({"error": "Flashcard not found."}), 404
    db.execute("DELETE FROM flashcards WHERE id = ?", (card_id,))
    db.commit()
    return jsonify({"message": "Flashcard deleted."})


@app.route("/api/flashcards/<int:card_id>/favorite", methods=["PATCH"])
def api_toggle_favorite(card_id):
    db = get_db()
    row = db.execute("SELECT favorite FROM flashcards WHERE id = ?", (card_id,)).fetchone()
    if not row:
        return jsonify({"error": "Flashcard not found."}), 404
    new_value = 0 if row["favorite"] else 1
    db.execute("UPDATE flashcards SET favorite=? WHERE id=?", (new_value, card_id))
    db.commit()
    return jsonify({"favorite": bool(new_value)})


@app.route("/api/subjects", methods=["GET"])
def api_subjects():
    rows = get_db().execute(
        "SELECT DISTINCT subject FROM flashcards WHERE subject IS NOT NULL"
    ).fetchall()
    return jsonify([r["subject"] for r in rows])


# ============================================================
# ROUTES — Quiz Mode
# ============================================================

@app.route("/api/quiz", methods=["GET"])
def api_quiz():
    """
    Builds a quiz from existing flashcards (MCQ / True-False / Fill-blank
    only — these have deterministic correct answers, so no extra AI calls
    are needed to grade them).
    """
    subject = request.args.get("subject", "").strip()
    count = min(request.args.get("count", 10, type=int) or 10, 30)

    placeholders = ",".join("?" * len(QUIZ_TYPES))
    query = f"SELECT * FROM flashcards WHERE card_type IN ({placeholders})"
    params = list(QUIZ_TYPES)
    if subject:
        query += " AND subject = ?"
        params.append(subject)

    rows = get_db().execute(query, params).fetchall()
    if not rows:
        return jsonify({"error": "No quiz-eligible flashcards found (need MCQ, True/False, or Fill-in-the-Blank cards)."}), 404

    cards = [row_to_dict(r) for r in rows]
    random.shuffle(cards)
    return jsonify(cards[:count])


# ============================================================
# ROUTES — Upload History
# ============================================================

@app.route("/api/documents", methods=["GET"])
def api_documents():
    rows = get_db().execute("SELECT * FROM documents ORDER BY created_at DESC LIMIT 50").fetchall()
    return jsonify([dict(r) for r in rows])


# ============================================================
# ROUTES — Export (CSV / JSON / PDF)
# ============================================================

@app.route("/api/export/csv", methods=["GET"])
def api_export_csv():
    rows = get_db().execute("SELECT * FROM flashcards ORDER BY created_at DESC").fetchall()
    buffer = StringIO()
    writer = csv.writer(buffer)
    writer.writerow(["Type", "Subject", "Topic", "Term/Question", "Answer", "Explanation", "Difficulty", "Favorite"])
    for r in rows:
        writer.writerow([r["card_type"], r["subject"] or "", r["topic"] or "",
                          r["term"] or r["question"] or "", r["answer"], r["explanation"] or "",
                          r["difficulty"], "Yes" if r["favorite"] else "No"])
    return Response(
        buffer.getvalue(),
        mimetype="text/csv",
        headers={"Content-Disposition": "attachment; filename=flashcards_export.csv"},
    )


@app.route("/api/export/json", methods=["GET"])
def api_export_json():
    rows = get_db().execute("SELECT * FROM flashcards ORDER BY created_at DESC").fetchall()
    payload = json.dumps([row_to_dict(r) for r in rows], indent=2)
    return Response(
        payload,
        mimetype="application/json",
        headers={"Content-Disposition": "attachment; filename=flashcards_export.json"},
    )


@app.route("/api/export/pdf", methods=["GET"])
def api_export_pdf():
    rows = get_db().execute("SELECT * FROM flashcards ORDER BY created_at DESC").fetchall()

    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleStyle", parent=styles["Title"], textColor=colors.HexColor("#4F46E5"))
    question_style = ParagraphStyle("Question", parent=styles["Heading3"], spaceBefore=10)
    answer_style = ParagraphStyle("Answer", parent=styles["BodyText"], textColor=colors.HexColor("#1F2937"))
    meta_style = ParagraphStyle("Meta", parent=styles["BodyText"], textColor=colors.grey, fontSize=8)

    story = [Paragraph("AI Flashcard Generator — Study Set", title_style), Spacer(1, 0.5 * cm)]
    for i, r in enumerate(rows, start=1):
        prompt_text = r["term"] or r["question"] or ""
        story.append(Paragraph(f"{i}. {prompt_text}", question_style))
        options = json.loads(r["options"]) if r["options"] else []
        for opt in options:
            story.append(Paragraph(f"&nbsp;&nbsp;&#8226; {opt}", answer_style))
        story.append(Paragraph(f"<b>Answer:</b> {r['answer']}", answer_style))
        if r["explanation"]:
            story.append(Paragraph(f"<i>{r['explanation']}</i>", answer_style))
        story.append(Paragraph(f"{r['card_type'].upper()} · {r['subject'] or 'General'} · {r['difficulty']}", meta_style))
        story.append(Spacer(1, 0.4 * cm))

    if not rows:
        story.append(Paragraph("No flashcards yet.", answer_style))

    doc.build(story)
    buffer.seek(0)
    return send_file(buffer, mimetype="application/pdf", as_attachment=True, download_name="flashcards_export.pdf")


# ============================================================
# ERROR HANDLERS — never fail silently
# ============================================================

@app.errorhandler(413)
def handle_too_large(e):
    return jsonify({"error": "Upload too large. Please reduce file size or number of files."}), 413


@app.errorhandler(404)
def handle_not_found(e):
    if request.path.startswith("/api/"):
        return jsonify({"error": "Resource not found."}), 404
    return send_from_directory(FRONTEND_DIR, "index.html")


@app.errorhandler(500)
def handle_server_error(e):
    logger.exception("Unhandled server error")
    return jsonify({"error": "Internal server error. Please try again."}), 500


# ============================================================
# ENTRYPOINT
# ============================================================

if __name__ == "__main__":
    init_db()
    port = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_DEBUG", "false").strip().lower() == "true"
    app.run(host="0.0.0.0", port=port, debug=debug)
else:
    init_db()
